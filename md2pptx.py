#!/usr/bin/env python

import mistune
from pptx import Presentation
from pprint import pprint
from parse import Parse,Page
from optparse import OptionParser
import warnings

def createpptx(mdfile,pptxfile,template):

    parse = Parse()
    parse.read(mdfile)
    title_name , allslides =  parse.get_title()

    prs = Presentation(template)

    xml_slides = prs.slides._sldIdLst
    removeslides = list(xml_slides)
    for slide in removeslides:
        xml_slides.remove(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[0])

    try:
        title = slide.shapes.title
        title.text =  title_name
    except:
        print("No Title in template file")
        exit()

    for slides in allslides:
        page = Page()
        page.split_elem(slides)
        sub_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(sub_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = page.title
        tf = body_shape.text_frame
        tf.text = page.text

    prs.save(pptxfile)


if __name__ == '__main__':

    # deal with warnings
    warnings.filterwarnings('ignore')

    parser = OptionParser()
    parser.add_option("-i", "--input", dest="input",
                      help="input markdown file", metavar="FILE")
    parser.add_option("-t", "--template", dest="template",
                      help="template(theme) pptx file", metavar="FILE")
    parser.add_option("-o", "--output", dest="output",
                      help="output markdown file", metavar="FILE")

    (options, args) = parser.parse_args()

    createpptx(options.input,options.output,options.template)
